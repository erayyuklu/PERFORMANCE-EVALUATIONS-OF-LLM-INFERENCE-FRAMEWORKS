#!/usr/bin/env python3
"""
run_gaia.py — Run GAIA benchmark against the LangGraph agent API.
=================================================================
Loads the GAIA validation set from HuggingFace, sends each question to
the agent, extracts the final answer, scores it with exact-match, and
writes per-question JSONL + summary JSON.

Usage:
    python run_gaia.py                           # defaults from config.env
    python run_gaia.py --levels 1 --limit 10     # Level 1 only, first 10
    python run_gaia.py --agent-url http://X:8000 # override agent URL
"""

import argparse
import asyncio
import json
import logging
import os
import re
import subprocess
import sys
import time
from datetime import datetime
from pathlib import Path

import httpx
import pandas as pd
from datasets import load_dataset
from huggingface_hub import snapshot_download

from scorer import score

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logging.basicConfig(
    stream=sys.stdout,
    level=logging.INFO,
    format="%(asctime)s  [%(levelname)s]  %(message)s",
)
logger = logging.getLogger(__name__)

SCRIPT_DIR = Path(__file__).resolve().parent

# ---------------------------------------------------------------------------
# Config — loaded from config.env then overridden by CLI args
# ---------------------------------------------------------------------------

def _load_env(path: Path) -> dict:
    """Parse a KEY=VALUE file, ignoring comments and blank lines."""
    env = {}
    if not path.exists():
        return env
    for line in path.read_text().splitlines():
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        key, _, value = line.partition("=")
        env[key.strip()] = value.strip().strip('"').strip("'")
    return env


_env = _load_env(SCRIPT_DIR / "config.env")

AGENT_K8S_NAMESPACE = _env.get("AGENT_K8S_NAMESPACE", "agent")
AGENT_K8S_SERVICE = _env.get("AGENT_K8S_SERVICE", "agent-service")
MODEL_FAMILY = _env.get("MODEL_FAMILY", "unknown-model")
QUANTIZATION = _env.get("QUANTIZATION", "noquantization")
AGENT_MODE = _env.get("AGENT_MODE", "single-agent")
GAIA_DATASET = _env.get("GAIA_DATASET", "gaia-benchmark/GAIA")
GAIA_CONFIG = _env.get("GAIA_CONFIG", "2023_all")
GAIA_SPLIT = _env.get("GAIA_SPLIT", "validation")
GAIA_LEVELS = _env.get("GAIA_LEVELS", "1,2,3")
LIMIT = int(_env.get("LIMIT", "0"))
OUTPUT_DIR = _env.get("OUTPUT_DIR", "../results")
REQUEST_TIMEOUT = int(_env.get("REQUEST_TIMEOUT", "180"))
MAX_CONCURRENT = int(_env.get("MAX_CONCURRENT", "4"))
MAX_RETRIES = int(_env.get("MAX_RETRIES", "3"))
RESUME = _env.get("RESUME", "false").lower() in ("true", "1", "yes")



def _discover_agent_url() -> str:
    """Discover the agent LoadBalancer URL via kubectl."""
    def _kubectl(jsonpath: str) -> str:
        return subprocess.check_output(
            ["kubectl", "get", "svc", AGENT_K8S_SERVICE, "-n", AGENT_K8S_NAMESPACE,
             "-o", f"jsonpath={jsonpath}"],
            text=True,
        ).strip()

    ip = _kubectl("{.status.loadBalancer.ingress[0].ip}")
    if not ip:
        raise RuntimeError(
            f"No external IP for {AGENT_K8S_SERVICE} in namespace {AGENT_K8S_NAMESPACE}. "
            f"Check: kubectl get svc {AGENT_K8S_SERVICE} -n {AGENT_K8S_NAMESPACE}"
        )
    port = _kubectl("{.spec.ports[0].port}")
    url = f"http://{ip}:{port}"
    logger.info(f"Discovered agent endpoint: {url}")
    return url

# ---------------------------------------------------------------------------
# GAIA prompt wrapper
# ---------------------------------------------------------------------------

GAIA_PROMPT_PREFIX = (
    "Answer the following question. Your response MUST end with exactly one "
    "line in this format:\n"
    "FINAL ANSWER: <your answer>\n"
    "where <your answer> is a short, precise response — a number, a name, "
    "a comma-separated list, or a short phrase, as the question requires. "
    "Do not include any explanation after the FINAL ANSWER line.\n\n"
)

# File-type extensions we can inline as plain text
TEXT_EXTENSIONS = {
    ".txt", ".csv", ".tsv", ".json", ".jsonl", ".md", ".py", ".js",
    ".html", ".xml", ".yaml", ".yml", ".log", ".sql", ".r", ".sh",
}

# Extensions we cannot extract text from
UNSUPPORTED_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".mp3", ".wav", ".mp4", ".mov"}


def _extract_file_content(file_path: str) -> str | None:
    """Extract text content from a file locally. Returns None if unsupported."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext in TEXT_EXTENSIONS:
        return Path(file_path).read_text(errors="replace")

    if ext == ".pdf":
        import pymupdf
        doc = pymupdf.open(file_path)
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n\n".join(pages).strip() or None

    if ext in (".xlsx", ".xls"):
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append("\t".join(str(c) if c is not None else "" for c in row))
            parts.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts).strip() or None

    if ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                paragraphs.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(paragraphs).strip() or None

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append("\t".join(cell.text for cell in row.cells))
            if texts:
                slides.append(f"=== Slide {i} ===\n" + "\n".join(texts))
        return "\n\n".join(slides).strip() or None

    return None


def _build_prompt(question: str, file_path: str | None) -> str:
    """Wrap a GAIA question with instructions and optional file content."""
    parts = [GAIA_PROMPT_PREFIX]

    if file_path and os.path.isfile(file_path):
        fname = os.path.basename(file_path)
        ext = os.path.splitext(file_path)[1].lower()

        if ext in UNSUPPORTED_EXTENSIONS:
            parts.append(
                f"[Attached file: {fname} — this is a {ext} file whose content "
                "cannot be provided as text.]\n\n"
            )
        else:
            try:
                content = _extract_file_content(file_path)
                if content:
                    parts.append(
                        f"The following file is attached: {fname}\n"
                        f"--- FILE CONTENT ---\n{content}\n--- END FILE ---\n\n"
                    )
                else:
                    parts.append(f"[Attached file: {fname} — extracted no content.]\n\n")
            except Exception as exc:
                parts.append(f"[Attached file: {fname} — could not read: {exc}]\n\n")

    parts.append(f"Question: {question}")
    return "".join(parts)


def _extract_final_answer(text: str) -> str:
    """Parse the 'FINAL ANSWER: ...' line from agent output."""
    # Search from the end since the answer should be the last such line
    matches = re.findall(r"FINAL ANSWER:\s*(.+)", text, re.IGNORECASE)
    if matches:
        return matches[-1].strip()
    # Fallback: return the last non-empty line
    lines = [l.strip() for l in text.strip().splitlines() if l.strip()]
    return lines[-1] if lines else ""


# ---------------------------------------------------------------------------
# Agent call
# ---------------------------------------------------------------------------

async def _call_agent(
    client: httpx.AsyncClient,
    task: str,
    semaphore: asyncio.Semaphore,
    agent_url: str,
    max_retries: int = 3,
    initial_delay: float = 2.0,
) -> tuple[str, float, dict]:
    """Send a task to the agent API and return (result_text, duration_ms, token_usage) with retries."""
    async with semaphore:
        attempt = 0
        while True:
            t0 = time.perf_counter()
            error_msg = None
            try:
                resp = await client.post(
                    f"{agent_url}/api/v1/agent/run",
                    json={"task": task},
                    timeout=REQUEST_TIMEOUT,
                )
                
                # Check status code
                if resp.status_code >= 400:
                    try:
                        err_body = resp.text
                    except Exception:
                        err_body = ""
                    raise httpx.HTTPStatusError(
                        f"Status {resp.status_code}: {err_body[:200]}",
                        request=resp.request,
                        response=resp,
                    )
                
                data = resp.json()
                result_text = data.get("result", "")
                
                # Check if result_text contains an error message or if predicted starts with [ERROR]
                predicted_check = _extract_final_answer(result_text)
                has_error = (
                    result_text.startswith("[ERROR]") or
                    predicted_check.startswith("[ERROR]") or
                    "http://status/500" in result_text or
                    "500 Internal Server Error" in result_text
                )
                
                if has_error:
                    raise ValueError(f"Agent returned error in result: {result_text[:200]}")
                
                duration_ms = (time.perf_counter() - t0) * 1000
                token_usage = {
                    "planner_input_tokens": data.get("planner_input_tokens", 0),
                    "planner_output_tokens": data.get("planner_output_tokens", 0),
                    "executor_input_tokens": data.get("executor_input_tokens", 0),
                    "executor_output_tokens": data.get("executor_output_tokens", 0),
                }
                return result_text, duration_ms, token_usage

            except Exception as exc:
                duration_ms = (time.perf_counter() - t0) * 1000
                attempt += 1
                error_msg = str(exc)
                logger.warning(
                    f"Agent call failed (attempt {attempt}/{max_retries + 1}): {exc} "
                    f"(duration: {duration_ms:.0f}ms)"
                )
                
                if attempt > max_retries:
                    return f"[ERROR] {error_msg}", duration_ms, {
                        "planner_input_tokens": 0,
                        "planner_output_tokens": 0,
                        "executor_input_tokens": 0,
                        "executor_output_tokens": 0,
                    }
                
                delay = initial_delay * (2 ** (attempt - 1))
                logger.info(f"Retrying in {delay:.1f}s...")
                await asyncio.sleep(delay)


# ---------------------------------------------------------------------------
# Main evaluation loop
# ---------------------------------------------------------------------------

async def run_evaluation(args: argparse.Namespace):
    agent_url = args.agent_url or _discover_agent_url()
    levels = [int(l) for l in (args.levels or GAIA_LEVELS).split(",")]
    limit = args.limit if args.limit is not None else LIMIT
    output_base = Path(SCRIPT_DIR / (args.output_dir or OUTPUT_DIR))
    model_family = args.model_family or MODEL_FAMILY
    quantization = args.quantization or QUANTIZATION
    agent_mode = args.agent_mode or AGENT_MODE
    max_concurrent = args.max_concurrent or MAX_CONCURRENT
    max_retries = args.max_retries if args.max_retries is not None else MAX_RETRIES
    resume = args.resume or RESUME

    # ── Download dataset ─────────────────────────────────────────────────
    logger.info(f"Downloading GAIA dataset: {GAIA_DATASET} ({GAIA_CONFIG}/{GAIA_SPLIT})")
    data_dir = snapshot_download(repo_id=GAIA_DATASET, repo_type="dataset")
    dataset = load_dataset(data_dir, GAIA_CONFIG, split=GAIA_SPLIT)
    logger.info(f"Loaded {len(dataset)} questions from GAIA {GAIA_SPLIT} split")

    # Filter by level
    dataset = dataset.filter(lambda x: int(x["Level"]) in levels)
    logger.info(f"After level filter ({levels}): {len(dataset)} questions")

    # Apply limit
    if limit > 0:
        dataset = dataset.select(range(min(limit, len(dataset))))
        logger.info(f"Limited to {len(dataset)} questions")

    # ── Prepare output directory ─────────────────────────────────────────
    # Convention: {OUTPUT_DIR}/{AGENT_MODE}/{MODEL_FAMILY}/{QUANTIZATION}/gaia
    run_dir = output_base / agent_mode / model_family / quantization / "gaia"
    run_dir.mkdir(parents=True, exist_ok=True)
    results_path = run_dir / "results.jsonl"
    summary_path = run_dir / "summary.json"

    logger.info(f"Results will be saved to: {run_dir}")
    logger.info(f"  Model family : {model_family}")
    logger.info(f"  Quantization : {quantization}")
    logger.info(f"  Agent mode   : {agent_mode}")
    logger.info(f"  Agent URL    : {agent_url}")
    logger.info(f"  Max concurrent: {max_concurrent}")
    logger.info("")

    # ── Run questions ────────────────────────────────────────────────────
    semaphore = asyncio.Semaphore(max_concurrent)
    
    # Load existing results if resuming
    existing_results = {}
    if resume and results_path.exists():
        logger.info(f"Found existing results file at {results_path}. Loading for resume...")
        try:
            with open(results_path, "r") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        record = json.loads(line)
                        tid = record.get("task_id")
                        if tid:
                            pred = record.get("predicted", "")
                            has_error = (
                                pred.startswith("[ERROR]") or
                                "http://status/500" in pred or
                                "500 Internal Server Error" in pred
                            )
                            if not has_error:
                                existing_results[tid] = record
                    except json.JSONDecodeError:
                        pass
            logger.info(f"Loaded {len(existing_results)} successful existing results.")
        except Exception as e:
            logger.warning(f"Failed to load existing results for resume: {e}")

    tasks_to_run = []
    final_results_placeholder = [None] * len(dataset)

    for idx, item in enumerate(dataset):
        task_id = item["task_id"]
        question = item["Question"]
        expected = item.get("Final answer", "")
        level = item["Level"]
        file_name = item.get("file_name", "") or ""

        # Check if we can resume this task
        if task_id in existing_results:
            final_results_placeholder[idx] = existing_results[task_id]
            continue

        # Resolve file attachment path
        file_path = None
        if file_name:
            candidate = os.path.join(data_dir, item.get("file_path", ""))
            if os.path.isfile(candidate):
                file_path = candidate

        prompt = _build_prompt(question, file_path)
        tasks_to_run.append((task_id, question, expected, level, file_name, prompt, idx))

    logger.info(
        f"Skipping {len(dataset) - len(tasks_to_run)} already completed tasks. "
        f"Running {len(tasks_to_run)} tasks."
    )

    results = []
    if tasks_to_run:
        async with httpx.AsyncClient() as client:
            async def _process(task_id, question, expected, level, file_name, prompt, idx):
                result_text, duration_ms, token_usage = await _call_agent(
                    client, prompt, semaphore, agent_url, max_retries=max_retries
                )
                predicted = _extract_final_answer(result_text)
                correct = score(predicted, expected)

                record = {
                    "task_id": task_id,
                    "level": level,
                    "question": question[:200],
                    "file_name": file_name,
                    "expected": expected,
                    "predicted": predicted,
                    "score": correct,
                    "duration_ms": round(duration_ms, 2),
                    **token_usage,
                }

                status = "CORRECT" if correct else "WRONG"
                logger.info(
                    f"[{status}] Level {level} | expected={expected!r} | "
                    f"predicted={predicted!r} | {duration_ms:.0f}ms"
                )
                return record, idx

            coros = [
                _process(tid, q, exp, lvl, fn, p, idx)
                for tid, q, exp, lvl, fn, p, idx in tasks_to_run
            ]
            new_results_with_idx = await asyncio.gather(*coros)
            for record, idx in new_results_with_idx:
                final_results_placeholder[idx] = record

    results = [r for r in final_results_placeholder if r is not None]

    # ── Write per-question results ───────────────────────────────────────
    with open(results_path, "w") as f:
        for r in results:
            f.write(json.dumps(r) + "\n")

    # ── Compute summary ──────────────────────────────────────────────────
    total = len(results)
    if total == 0:
        logger.warning("No questions were evaluated. Check your --levels and --limit flags.")
        summary = {"total": 0, "correct": 0, "accuracy": 0.0}
        with open(summary_path, "w") as f:
            json.dump(summary, f, indent=2)
        return

    df = pd.DataFrame(results)
    correct = int(df["score"].sum())
    accuracy = correct / total if total > 0 else 0.0

    # Calculate token statistics
    total_planner_input_tokens = int(df["planner_input_tokens"].sum()) if "planner_input_tokens" in df else 0
    total_planner_output_tokens = int(df["planner_output_tokens"].sum()) if "planner_output_tokens" in df else 0
    total_executor_input_tokens = int(df["executor_input_tokens"].sum()) if "executor_input_tokens" in df else 0
    total_executor_output_tokens = int(df["executor_output_tokens"].sum()) if "executor_output_tokens" in df else 0

    avg_planner_input_tokens = round(float(df["planner_input_tokens"].mean()), 2) if "planner_input_tokens" in df and total > 0 else 0.0
    avg_planner_output_tokens = round(float(df["planner_output_tokens"].mean()), 2) if "planner_output_tokens" in df and total > 0 else 0.0
    avg_executor_input_tokens = round(float(df["executor_input_tokens"].mean()), 2) if "executor_input_tokens" in df and total > 0 else 0.0
    avg_executor_output_tokens = round(float(df["executor_output_tokens"].mean()), 2) if "executor_output_tokens" in df and total > 0 else 0.0

    per_level = {}
    for lvl in sorted(df["level"].unique()):
        lvl_df = df[df["level"] == lvl]
        lvl_total = len(lvl_df)
        lvl_correct = int(lvl_df["score"].sum())
        per_level[f"level_{lvl}"] = {
            "total": lvl_total,
            "correct": lvl_correct,
            "accuracy": round(lvl_correct / lvl_total, 4) if lvl_total > 0 else 0.0,
        }

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    summary = {
        "timestamp": timestamp,
        "model_family": model_family,
        "quantization": quantization,
        "agent_mode": agent_mode,
        "agent_url": agent_url,
        "dataset": GAIA_DATASET,
        "config": GAIA_CONFIG,
        "split": GAIA_SPLIT,
        "levels": levels,
        "total": total,
        "correct": correct,
        "accuracy": round(accuracy, 4),
        "per_level": per_level,
        "avg_duration_ms": round(float(df["duration_ms"].mean()), 2) if total > 0 else 0.0,
        "total_planner_input_tokens": total_planner_input_tokens,
        "total_planner_output_tokens": total_planner_output_tokens,
        "total_executor_input_tokens": total_executor_input_tokens,
        "total_executor_output_tokens": total_executor_output_tokens,
        "avg_planner_input_tokens": avg_planner_input_tokens,
        "avg_planner_output_tokens": avg_planner_output_tokens,
        "avg_executor_input_tokens": avg_executor_input_tokens,
        "avg_executor_output_tokens": avg_executor_output_tokens,
    }

    with open(summary_path, "w") as f:
        json.dump(summary, f, indent=2)

    # ── Print summary ────────────────────────────────────────────────────
    logger.info("")
    logger.info("=" * 60)
    logger.info("GAIA Evaluation Summary")
    logger.info("=" * 60)
    logger.info(f"Overall: {correct}/{total} = {accuracy:.1%}")
    for lvl_key, lvl_data in per_level.items():
        logger.info(
            f"  {lvl_key}: {lvl_data['correct']}/{lvl_data['total']} "
            f"= {lvl_data['accuracy']:.1%}"
        )
    logger.info(f"Avg duration: {summary['avg_duration_ms']:.0f}ms")
    logger.info(f"Planner Tokens - Total Input: {total_planner_input_tokens} | Total Output: {total_planner_output_tokens}")
    logger.info(f"Executor Tokens - Total Input: {total_executor_input_tokens} | Total Output: {total_executor_output_tokens}")
    logger.info(f"Planner Tokens - Avg Input: {avg_planner_input_tokens:.1f} | Avg Output: {avg_planner_output_tokens:.1f}")
    logger.info(f"Executor Tokens - Avg Input: {avg_executor_input_tokens:.1f} | Avg Output: {avg_executor_output_tokens:.1f}")
    logger.info(f"Results saved to: {run_dir}")
    logger.info("=" * 60)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Run GAIA benchmark against the agent API.")
    parser.add_argument("--agent-url", type=str, default=None, help="Agent base URL")
    parser.add_argument("--levels", type=str, default=None, help="Comma-separated levels (e.g. 1,2)")
    parser.add_argument("--limit", type=int, default=None, help="Max questions to run (0=all)")
    parser.add_argument("--output-dir", type=str, default=None, help="Output base directory")
    parser.add_argument("--model-family", type=str, default=None, help="Model family (e.g. qwen-3)")
    parser.add_argument("--quantization", type=str, default=None, help="Quantization config")
    parser.add_argument("--agent-mode", type=str, default=None, help="Agent mode (single-agent|planner-only|planner-executor)")
    parser.add_argument("--max-concurrent", type=int, default=None, help="Max concurrent requests")
    parser.add_argument("--max-retries", type=int, default=None, help="Max retries for transient errors")
    parser.add_argument("--resume", action="store_true", help="Resume from existing results.jsonl (skipping successful tasks)")
    args = parser.parse_args()
    asyncio.run(run_evaluation(args))


if __name__ == "__main__":
    main()